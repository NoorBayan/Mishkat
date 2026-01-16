import os
import shutil
import subprocess
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# Imports from src
from ..utils.text import is_arabic, estimate_font_size
from ..config import (
    ARABIC_FONT_MAIN,
    ARABIC_FONT_META,
    LATIN_FONT_MAIN,
    LATIN_FONT_META
)

# =====================================================
# Main Function
# =====================================================

def render_quran_slide_from_template(
    template_pptx: str,
    output_pptx: str,
    ayah_number: int,
    surah_name: str,
    ayah_text: str
):
    """
    Renders a Quran slide by copying a BASE template and writing content
    into the copied presentation (template-safe).
    """

    # -------------------------------------------------
    # 1) Copy template → output
    # -------------------------------------------------
    os.makedirs(os.path.dirname(output_pptx), exist_ok=True)
    shutil.copy(template_pptx, output_pptx)

    # -------------------------------------------------
    # 2) Open copied presentation
    # -------------------------------------------------
    prs = Presentation(output_pptx)
    slide = prs.slides[0]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    arabic = is_arabic(ayah_text)

    # -------------------------------------------------
    # 3) Remove ONLY text boxes (safe cleanup)
    # -------------------------------------------------
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            slide.shapes._spTree.remove(shape._element)

    # -------------------------------------------------
    # 4) Layout configuration
    # -------------------------------------------------
    box_left = slide_width * 0.08
    box_top = slide_height * 0.28
    box_width = slide_width * 0.84
    box_height = slide_height * 0.34

    font_size = estimate_font_size(
        text=ayah_text,
        box_width_ratio=0.84,
        box_height_ratio=0.34,
        slide_width=slide_width,
        slide_height=slide_height,
        is_arabic_text=arabic
    )

    # -------------------------------------------------
    # 5) Main Ayah Text
    # -------------------------------------------------
    main_box = slide.shapes.add_textbox(
        left=box_left,
        top=box_top,
        width=box_width,
        height=box_height
    )

    tf = main_box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.rtl = arabic
    p.text = f"﴿ {ayah_text} ﴾" if arabic else ayah_text

    run = p.runs[0]
    run.font.name = ARABIC_FONT_MAIN if arabic else LATIN_FONT_MAIN

    run.font.size = Pt(font_size)

    # -------------------------------------------------
    # 6) Footer
    # -------------------------------------------------
    footer_box = slide.shapes.add_textbox(
        left=slide_width * 0.25,
        top=slide_height * 0.70,
        width=slide_width * 0.50,
        height=slide_height * 0.10
    )

    footer_tf = footer_box.text_frame
    footer_tf.clear()
    footer_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p_footer = footer_tf.paragraphs[0]
    p_footer.alignment = PP_ALIGN.CENTER
    p_footer.rtl = arabic
    p_footer.text = (
        f"سورة {surah_name} • الآية {ayah_number}"
        if arabic
        else f"Surah {surah_name} – Ayah {ayah_number}"
    )

    run_footer = p_footer.runs[0]
    run_footer.font.name = ARABIC_FONT_META if arabic else LATIN_FONT_META

    run_footer.font.size = Pt(18)

    # -------------------------------------------------
    # 7) Save rendered slide
    # -------------------------------------------------
    prs.save(output_pptx)


def pptx_to_png(pptx_path: str, output_dir: str):
    os.makedirs(output_dir, exist_ok=True)

    cmd = [
        "libreoffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to", "png",
        "--outdir", output_dir,
        pptx_path
    ]

    subprocess.run(cmd, check=True)

    return output_dir
